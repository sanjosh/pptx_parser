

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def get_fmt(s):
    return f"[top={s.top}, h={s.height} left={s.left}, w={s.width}]"


def print_shape(shape):
    """
    https://python-pptx.readthedocs.io/en/latest/api/enum/MsoShapeType.html
    :param shape:
    :return:
    """
    if shape.shape_type == MSO_SHAPE_TYPE.LINE:
        line = shape.line
        print(f"line : {line.dash_style} {get_fmt(shape)}")
    elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        print(f'text={shape.text} {get_fmt(shape)}')
    elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        # {shape.auto_shape_type} fails on line
        # https://python-pptx.readthedocs.io/en/latest/api/enum/MsoAutoShapeType.html
        print(f'autoshape: {shape.name} {shape.text}')
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for shape in shape.shapes:
            print_shape(shape)
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        print(f'pic {shape.shape_type}')
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        print(f'table {shape.shape_type}')
    else:
        print(f'todo {shape.shape_type}')


def print_hi():
    path_to_presentation = './xyz_gb_prodg.pptx'
    prs = Presentation(path_to_presentation)
    print(f"{prs.slide_width} and {prs.slide_height}")
    for idx, slide in enumerate(prs.slides):
        print(f"{idx}:SLIDE NAME = {slide.slide_layout.name}")
        if idx != 18:
            continue
        print(f"page={idx}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                print(f'text: {shape.text} {get_fmt(shape)}')
                # for paragraph in shape.text_frame.paragraphs:
                #     for run in paragraph.runs:
                #         text_runs.append(run.text)
            elif shape.has_chart:
                print('chart')
                # print(shape.chart)
            elif shape.has_table:
                print('table')
                # print(shape.table)
            else:
                print(f"name={shape.name}, id={shape.shape_id}, type={shape.shape_type}, obj={shape}")
                print_shape(shape)



    # print(text_runs)

# Press the green button in the gutter to run the script.
if __name__ == '__main__':
    print_hi()

# See PyCharm help at https://www.jetbrains.com/help/pycharm/
